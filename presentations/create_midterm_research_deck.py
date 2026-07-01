from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from create_midterm_findings_deck import (
    COLORS,
    DATA,
    FIG,
    ROOT,
    WIDE_H,
    WIDE_W,
    add_bullets,
    add_card,
    add_footer,
    add_metric,
    add_picture_fit,
    add_textbox,
    add_title,
    blank,
    fmt_pct,
    fmt_rank,
    read_csv,
)


OUT = ROOT / "presentations" / "midterm_research_presentation_v2.pptx"


def build_deck() -> tuple[Path, int]:
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H

    elec_rank = {r["technology"]: r for r in read_csv(DATA / "2026-06-23-NPV_Ranking_Electricity_summary.csv")}
    cement_rank = {r["technology"]: r for r in read_csv(DATA / "2026-06-23-NPV_Ranking_Cement_summary.csv")}
    macc_det = {r["technology"]: r for r in read_csv(DATA / "2026-06-26-Cement_MACC_Deterministic.csv")}
    macc_sim = {r["technology"]: r for r in read_csv(DATA / "2026-06-26-Cement_MACC_Simulated.csv")}

    slides = []

    # 1
    s = blank(prs)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = COLORS["light"]
    add_textbox(s, "Midterm Research Presentation", Inches(0.75), Inches(0.95), Inches(9.6), Inches(0.55), 34, True, "blue")
    add_textbox(
        s,
        "Impact of technology-development uncertainty on carbon-capture demand uncertainty",
        Inches(0.78),
        Inches(1.62),
        Inches(10.8),
        Inches(0.45),
        21,
        True,
    )
    add_textbox(
        s,
        "Research midterm | Master thesis | ETH Zurich | July 2026",
        Inches(0.8),
        Inches(2.22),
        Inches(8.8),
        Inches(0.3),
        13,
        False,
        "muted",
    )
    add_card(
        s,
        "Central idea",
        "Carbon capture demand is modeled as an uncertain outcome of competition between CCS and alternative decarbonization technologies.",
        Inches(0.85),
        Inches(4.25),
        Inches(5.75),
        Inches(1.25),
        "green",
        15,
        13,
    )
    add_card(
        s,
        "Current midterm scope",
        "The implemented model currently compares electricity and cement technologies using deterministic and Monte Carlo NPV analysis.",
        Inches(6.9),
        Inches(4.25),
        Inches(5.45),
        Inches(1.25),
        "blue",
        15,
        13,
    )
    slides.append(s)

    # 2
    s = blank(prs)
    add_title(s, "Research Question", "The project question from the thesis proposal")
    add_card(
        s,
        "Main research question",
        "Which emitting sectors should implement carbon capture, and which should investigate alternative decarbonization technologies?",
        Inches(1.0),
        Inches(1.55),
        Inches(11.3),
        Inches(1.25),
        "blue",
        18,
        16,
    )
    add_bullets(
        s,
        [
            "The question is sector-specific: CCS may be attractive in one application and unattractive in another.",
            "The question is probabilistic: technology costs, operating costs, energy requirements, energy prices, and carbon prices are uncertain.",
            "The question is competitive: CCS must be compared against non-capture options, not evaluated in isolation.",
        ],
        Inches(1.05),
        Inches(3.35),
        Inches(10.9),
        Inches(2.0),
        18,
    )
    add_footer(s, 2)
    slides.append(s)

    # 3
    s = blank(prs)
    add_title(s, "What We Want To Find Out", "Research objective translated into measurable outputs")
    add_card(
        s,
        "1. Economic preference",
        "Which technology has the highest NPV under a given sampled future?",
        Inches(0.8),
        Inches(1.55),
        Inches(3.65),
        Inches(1.25),
        "blue",
    )
    add_card(
        s,
        "2. CCS probability",
        "How often does CCS rank first or near the top across uncertain techno-economic conditions?",
        Inches(4.85),
        Inches(1.55),
        Inches(3.65),
        Inches(1.25),
        "green",
    )
    add_card(
        s,
        "3. Abatement value",
        "How much emissions reduction does each option provide, and at what marginal cost?",
        Inches(8.9),
        Inches(1.55),
        Inches(3.65),
        Inches(1.25),
        "amber",
    )
    add_bullets(
        s,
        [
            "The midterm model therefore reports NPV distributions, ranking probabilities, sensitivity results, and cement marginal abatement costs.",
            "The final thesis should connect these outputs to sector-level CCS demand: strong, conditional, weak, or not yet assessable.",
        ],
        Inches(1.0),
        Inches(3.55),
        Inches(11.1),
        Inches(1.35),
        18,
    )
    add_footer(s, 3)
    slides.append(s)

    # 4
    s = blank(prs)
    add_title(s, "Experiment Setup", "Current modeled sectors and technology options")
    add_card(
        s,
        "Electricity sector",
        "Hard coal, hard coal + CCS, CCGT, CCGT + CCS, nuclear, offshore wind, onshore wind, PV, and biogas.",
        Inches(0.8),
        Inches(1.55),
        Inches(5.8),
        Inches(1.35),
        "blue",
    )
    add_card(
        s,
        "Cement sector",
        "BAU, CCS, clinker substitution, alternative fuels, efficiency improvement, waste heat recovery, process heat integration, electrification, and electrolysis.",
        Inches(6.9),
        Inches(1.55),
        Inches(5.8),
        Inches(1.35),
        "green",
    )
    add_bullets(
        s,
        [
            "Within each sector, technologies are normalized to the same annual output so rankings are not driven by plant size.",
            "Electricity technologies are modeled as absolute production options.",
            "Cement contains both absolute technologies and retrofit technologies relative to BAU.",
        ],
        Inches(1.0),
        Inches(3.55),
        Inches(11.1),
        Inches(1.5),
        17,
    )
    add_footer(s, 4)
    slides.append(s)

    # 5
    s = blank(prs)
    add_title(s, "Model Assumptions", "Baseline setup used for the current midterm results")
    metrics = [
        ("1 TWh/y", "electricity output", "blue"),
        ("1 Mt/y", "cement output", "green"),
        ("80", "EUR/tCO2 carbon price", "amber"),
        ("8%", "discount rate", "blue"),
    ]
    for i, (value, label, color) in enumerate(metrics):
        add_metric(s, value, label, Inches(0.85 + i * 3.05), Inches(1.55), Inches(2.45), color)
    add_card(
        s,
        "Revenue assumptions",
        "Electricity revenue uses 94.07 EUR/MWh. Cement revenue uses 150 EUR/t.",
        Inches(0.85),
        Inches(3.2),
        Inches(5.6),
        Inches(1.05),
        "blue",
    )
    add_card(
        s,
        "Uncertainty approach",
        "Monte Carlo simulation samples technology and market inputs. Default setup uses 100,000 draws and seed 42.",
        Inches(6.9),
        Inches(3.2),
        Inches(5.6),
        Inches(1.05),
        "green",
    )
    add_card(
        s,
        "Decision metric",
        "Higher NPV is better. Results are evaluated as total NPV and normalized NPV: EUR/MWh for electricity and EUR/t for cement.",
        Inches(0.85),
        Inches(4.75),
        Inches(11.65),
        Inches(0.95),
        "amber",
        14,
        12,
    )
    add_footer(s, 5)
    slides.append(s)

    # 6
    s = blank(prs)
    add_title(s, "Input Parameters", "Main parameter groups driving the experiment")
    add_bullets(
        s,
        [
            "Technology cost: CAPEX, fixed OPEX, variable OPEX.",
            "Energy requirements: fuel consumption, electricity consumption, full-load hours.",
            "Environmental performance: direct emissions and residual emissions after CCS.",
            "Market and policy inputs: electricity price, cement price, fuel prices, carbon price, discount rate.",
            "Cement retrofit inputs: CAPEX/OPEX changes and reductions in fuel use, electricity use, and emissions relative to BAU.",
        ],
        Inches(0.9),
        Inches(1.55),
        Inches(6.2),
        Inches(3.4),
        17,
    )
    add_card(
        s,
        "Recent modeling correction",
        "Alternative fuels in cement now use a 25-60% alternative-fuel share with a blended fuel price. This avoids assuming 100% alternative-fuel use while emissions reductions only reflect partial substitution.",
        Inches(7.35),
        Inches(1.7),
        Inches(4.75),
        Inches(1.65),
        "amber",
        15,
        12,
    )
    add_card(
        s,
        "Research implication",
        "Model assumptions directly influence CCS demand conclusions, so transparent parameter ownership is part of the research method.",
        Inches(7.35),
        Inches(3.75),
        Inches(4.75),
        Inches(1.35),
        "green",
        15,
        12,
    )
    add_footer(s, 6)
    slides.append(s)

    # 7
    s = blank(prs)
    add_title(s, "Result 1: Electricity", "Renewables dominate current NPV rankings")
    add_picture_fit(s, FIG / "2026-06-23-Average_NPV_Rank_Electricity.png", Inches(0.65), Inches(1.42), Inches(7.15), Inches(4.72))
    add_card(
        s,
        "Ranking result",
        f"PV has average rank {fmt_rank(elec_rank['pv']['average_rank'])} and rank-1 probability {fmt_pct(elec_rank['pv']['probability_rank_1'])}. "
        f"Onshore and offshore wind have top-3 probabilities of {fmt_pct(elec_rank['wind_onshore']['probability_top_3'])} and {fmt_pct(elec_rank['wind_offshore']['probability_top_3'])}.",
        Inches(8.05),
        Inches(1.55),
        Inches(4.55),
        Inches(1.8),
        "green",
        15,
        12,
    )
    add_card(
        s,
        "CCS signal",
        f"CCGT+CCS top-3 probability is {fmt_pct(elec_rank['ccgt_ccs']['probability_top_3'])}; hard coal+CCS top-3 probability is {fmt_pct(elec_rank['hard_coal_ccs']['probability_top_3'])}.",
        Inches(8.05),
        Inches(3.75),
        Inches(4.55),
        Inches(1.25),
        "blue",
        15,
        12,
    )
    add_footer(s, 7)
    slides.append(s)

    # 8
    s = blank(prs)
    add_title(s, "Result 2: Cement", "CCS is competitive, but not uncontested")
    add_picture_fit(s, FIG / "2026-06-23-Average_NPV_Rank_Cement.png", Inches(0.65), Inches(1.42), Inches(7.15), Inches(4.72))
    add_card(
        s,
        "CCS result",
        f"CCS has the highest rank-1 probability in cement: {fmt_pct(cement_rank['ccs']['probability_rank_1'])}. "
        f"Its average rank is {fmt_rank(cement_rank['ccs']['average_rank'])}, showing that outcomes are highly scenario-dependent.",
        Inches(8.05),
        Inches(1.55),
        Inches(4.55),
        Inches(1.65),
        "green",
        15,
        12,
    )
    add_card(
        s,
        "Non-capture competition",
        "Process heat integration, waste heat recovery, and clinker substitution are strong competitors because they improve economics without very high CAPEX.",
        Inches(8.05),
        Inches(3.6),
        Inches(4.55),
        Inches(1.35),
        "blue",
        15,
        12,
    )
    add_footer(s, 8)
    slides.append(s)

    # 9
    s = blank(prs)
    add_title(s, "Result 3: Cement MACC", "CCS has large abatement potential at moderate cost")
    add_picture_fit(s, FIG / "2026-06-26-Cement_MACC_Deterministic.png", Inches(0.6), Inches(1.38), Inches(7.75), Inches(4.6))
    add_card(
        s,
        "Deterministic abatement cost",
        f"Efficiency improvement: {float(macc_det['efficiency_improvement']['abatement_cost_eur_per_tco2']):.1f} EUR/tCO2.\n"
        f"Process heat integration: {float(macc_det['process_heat_integration']['abatement_cost_eur_per_tco2']):.1f} EUR/tCO2.\n"
        f"CCS: {float(macc_det['ccs']['abatement_cost_eur_per_tco2']):.1f} EUR/tCO2.",
        Inches(8.65),
        Inches(1.55),
        Inches(3.8),
        Inches(1.7),
        "green",
        15,
        12,
    )
    add_card(
        s,
        "Simulated CCS range",
        f"CCS median: {float(macc_sim['ccs']['abatement_cost_median_eur_per_tco2']):.1f} EUR/tCO2.\n"
        f"p05-p95: {float(macc_sim['ccs']['abatement_cost_p05_eur_per_tco2']):.1f} to {float(macc_sim['ccs']['abatement_cost_p95_eur_per_tco2']):.1f} EUR/tCO2.",
        Inches(8.65),
        Inches(3.65),
        Inches(3.8),
        Inches(1.35),
        "blue",
        15,
        12,
    )
    add_footer(s, 9)
    slides.append(s)

    # 10
    s = blank(prs)
    add_title(s, "Result 4: Sensitivity", "Which assumptions drive the results?")
    add_picture_fit(s, FIG / "2026-06-24-Sensitivity_Heatmap_Standardized_Cement.png", Inches(0.55), Inches(1.38), Inches(5.95), Inches(4.75))
    add_picture_fit(s, FIG / "2026-06-24-Sensitivity_Heatmap_Standardized_Electricity.png", Inches(6.75), Inches(1.38), Inches(5.95), Inches(4.75))
    add_footer(s, 10)
    slides.append(s)

    # 11
    s = blank(prs)
    add_title(s, "Discussion", "Interpreting the midterm evidence")
    add_bullets(
        s,
        [
            "Electricity: current model results suggest weak CCS demand because renewable generation dominates economically.",
            "Cement: current model results suggest conditional CCS demand because CCS frequently ranks first and offers deep abatement, but low-cost retrofits compete strongly.",
            "The same technology, CCS, has different demand implications by sector because alternatives, emissions sources, and energy requirements differ.",
            "This supports the thesis framing: CCS demand uncertainty should be modeled through sector-specific competition under uncertainty.",
        ],
        Inches(0.9),
        Inches(1.55),
        Inches(7.0),
        Inches(3.9),
        18,
    )
    add_card(
        s,
        "Preliminary conclusion",
        "CCS should not be treated as a universal decarbonization default. Its role is strongest where alternatives are costly, incomplete, or unable to address process emissions.",
        Inches(8.35),
        Inches(2.0),
        Inches(3.9),
        Inches(1.8),
        "green",
        15,
        13,
    )
    add_footer(s, 11)
    slides.append(s)

    # 12
    s = blank(prs)
    add_title(s, "Quality-Level Improvements", "How the research design can become stronger")
    add_bullets(
        s,
        [
            "Improve source quality and traceability for every parameter range, especially emerging technologies.",
            "Represent uncertainty correlations where independent sampling is physically weak, for example fuel use and emissions.",
            "Separate market, policy, and technology-development uncertainty more explicitly.",
            "Use scenario families to test robustness rather than relying only on one central assumption set.",
            "Clarify the decision criterion: highest expected NPV, probability of rank 1, probability of top 3, or abatement-adjusted value.",
            "Move from technology ranking to sector-level CCS demand categories.",
        ],
        Inches(0.9),
        Inches(1.55),
        Inches(11.1),
        Inches(4.0),
        17,
    )
    add_footer(s, 12)
    slides.append(s)

    # 13
    s = blank(prs)
    add_title(s, "Outlook", "How the project proceeds scientifically")
    add_card(
        s,
        "Broaden sector coverage",
        "Next modeled sectors: aluminium and chemical industries. These extend the thesis beyond electricity and cement toward industrial CCS demand.",
        Inches(0.8),
        Inches(1.55),
        Inches(5.65),
        Inches(1.5),
        "blue",
    )
    add_card(
        s,
        "Deepen interpretation",
        "For each sector, classify CCS demand as robust, conditional, weak, or dependent on unresolved assumptions.",
        Inches(6.85),
        Inches(1.55),
        Inches(5.65),
        Inches(1.5),
        "green",
    )
    add_card(
        s,
        "Compare across sectors",
        "Identify common drivers of CCS competitiveness: process emissions, energy penalties, electricity intensity, carbon price exposure, and retrofit availability.",
        Inches(0.8),
        Inches(3.55),
        Inches(5.65),
        Inches(1.5),
        "amber",
    )
    add_card(
        s,
        "Build final thesis argument",
        "Translate probabilistic technology rankings into a clear narrative about where carbon capture demand is likely, uncertain, or displaced by alternatives.",
        Inches(6.85),
        Inches(3.55),
        Inches(5.65),
        Inches(1.5),
        "blue",
    )
    add_footer(s, 13)
    slides.append(s)

    # 14
    s = blank(prs)
    add_title(s, "Technical Steps TODO", "Implementation tasks separated from the research outlook")
    add_bullets(
        s,
        [
            "Regenerate cement summary figures and CSV outputs after the alternative-fuels correction.",
            "Add automated identity checks for NPV, fuel costs, emissions costs, and normalized units.",
            "Create reusable templates for adding aluminium and chemical-sector technologies.",
            "Document parameter sources directly alongside assumptions.",
            "Keep generated outputs reproducible with sample size, seed, date, and model version.",
        ],
        Inches(0.9),
        Inches(1.55),
        Inches(11.0),
        Inches(3.3),
        18,
    )
    add_footer(s, 14)
    slides.append(s)

    # 15
    s = blank(prs)
    add_title(s, "Appendix: Baseline NPV Figures", "Representative deterministic and mean results")
    add_picture_fit(s, FIG / "2026-06-23-Mean_NPV_per_MWh_Electricity.png", Inches(0.75), Inches(1.48), Inches(5.65), Inches(4.65))
    add_picture_fit(s, FIG / "2026-06-23-Mean_NPV_per_t_Cement.png", Inches(6.85), Inches(1.48), Inches(5.65), Inches(4.65))
    add_footer(s, 15)
    slides.append(s)

    prs.save(OUT)
    return OUT, len(slides)


if __name__ == "__main__":
    path, count = build_deck()
    print(f"Wrote {path} with {count} slides")
