from __future__ import annotations

import sys
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))
if str(Path(__file__).resolve().parent) not in sys.path:
    sys.path.insert(0, str(Path(__file__).resolve().parent))

from cement.cement_macc import deterministic_cement_macc, plot_cement_macc, simulated_cement_macc
from cement.cement_npv_monte_carlo import DEFAULT_RETROFIT_BAU_MODE, simulate_cement_results
from cement.cement_npv_summary_figures import CEMENT_TECHNOLOGY_LABELS, calculate_cement_npv_rankings_from_results
from electricity.electricity_npv_monte_carlo import simulate_electricity_results
from electricity.electricity_npv_summary_figures import ELECTRICITY_TECHNOLOGY_LABELS, calculate_electricity_npv_rankings_from_results
from general_parameters import CARBON_PRICE_EUR_PER_T, INTEREST_RATE
from npv_finance import calculate_npv
from sensitivity_analysis import METRIC_SPECIFIC
from sensitivity_deep_dive import build_sensitivity_heatmap_figure, standardized_sensitivity

from create_midterm_findings_deck import (
    COLORS,
    WIDE_H,
    WIDE_W,
    add_bullets,
    add_card,
    add_footer,
    add_metric,
    add_picture_fit,
    add_textbox,
    add_title,
    blank,
)
from pptx import Presentation
from pptx.util import Inches


OUT = ROOT / "presentations" / "midterm_deep_dive_v3.pptx"
REPORT = ROOT / "presentations" / "midterm_deep_dive_findings.md"
ASSET_DIR = ROOT / "presentations" / "deep_dive_assets"
SAMPLE_SIZE = 100_000
RANDOM_SEED = 42


def label_map(sector: str) -> dict[str, str]:
    return CEMENT_TECHNOLOGY_LABELS if sector == "cement" else ELECTRICITY_TECHNOLOGY_LABELS


def summarize(results: dict, labels: dict[str, str], specific_col: str) -> pd.DataFrame:
    rows = []
    for tech, result in results.items():
        specific = np.asarray(result[specific_col], dtype=float)
        total = np.asarray(result["npv_eur"], dtype=float) / 1e6
        rows.append(
            {
                "technology": tech,
                "label": labels.get(tech, tech),
                "mean_specific": float(np.mean(specific)),
                "median_specific": float(np.median(specific)),
                "p05_specific": float(np.percentile(specific, 5)),
                "p95_specific": float(np.percentile(specific, 95)),
                "mean_meur": float(np.mean(total)),
                "non_negative_share": float(np.mean(total >= 0)),
            }
        )
    return pd.DataFrame(rows).sort_values("mean_specific", ascending=False)


def scenario_npv(results: dict, carbon_price: float, discount_rate: float) -> np.ndarray:
    initial_capex = np.asarray(results["initial_capex_eur"], dtype=float)
    baseline_cash_flow = np.asarray(results["annual_net_cash_flow_eur"], dtype=float)
    baseline_emissions_cost = np.asarray(results["annual_emissions_cost_eur"], dtype=float)
    baseline_carbon_price = np.asarray(results["carbon_price_eur_per_t"], dtype=float)
    annual_emissions = np.divide(
        baseline_emissions_cost,
        baseline_carbon_price,
        out=np.zeros_like(baseline_emissions_cost),
        where=baseline_carbon_price != 0,
    )
    scenario_cash_flow = (
        baseline_cash_flow + baseline_emissions_cost - annual_emissions * carbon_price
    )
    lifetime_years = int(np.asarray(results["lifetime_years"])[0])
    return calculate_npv(
        initial_capex_eur=initial_capex,
        annual_net_cash_flow_eur=scenario_cash_flow,
        lifetime_years=lifetime_years,
        discount_rate=discount_rate,
    )


def build_scenario_table(
    sector: str,
    results: dict,
    labels: dict[str, str],
    denominator_col: str,
    kind: str,
) -> pd.DataFrame:
    values = (
        {"Low": 40.0, "Medium": CARBON_PRICE_EUR_PER_T.value, "High": 120.0}
        if kind == "co2"
        else {"Low": 0.04, "Medium": INTEREST_RATE.value, "High": 0.12}
    )
    rows = []
    for tech, result in results.items():
        denominator = np.asarray(result[denominator_col], dtype=float)
        for case, value in values.items():
            carbon_price = value if kind == "co2" else CARBON_PRICE_EUR_PER_T.value
            discount_rate = value if kind == "discount" else INTEREST_RATE.value
            specific = scenario_npv(result, carbon_price, discount_rate) / denominator
            rows.append(
                {
                    "sector": sector,
                    "technology": tech,
                    "label": labels.get(tech, tech),
                    "scenario": case,
                    "value": value,
                    "mean_specific": float(np.mean(specific)),
                    "non_negative_share": float(np.mean(specific >= 0)),
                }
            )
    return pd.DataFrame(rows)


def save_bar_chart(df: pd.DataFrame, path: Path, title: str, xlabel: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    plot_df = df.sort_values("mean_specific")
    colors = ["#2f6f4e" if value >= 0 else "#9b4a4a" for value in plot_df["mean_specific"]]
    fig, ax = plt.subplots(figsize=(8.5, 4.6), dpi=180)
    ax.barh(plot_df["label"], plot_df["mean_specific"], color=colors)
    ax.axvline(0, color="#222222", linewidth=0.8)
    ax.set_title(title, loc="left", fontsize=12, fontweight="bold")
    ax.set_xlabel(xlabel)
    ax.grid(axis="x", color="#e5e7eb", linestyle=(0, (2, 4)))
    for spine in ("top", "right", "left"):
        ax.spines[spine].set_visible(False)
    ax.tick_params(axis="y", length=0, labelsize=8)
    ax.tick_params(axis="x", labelsize=8)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def save_rank_chart(df: pd.DataFrame, labels: dict[str, str], path: Path, title: str) -> None:
    plot_df = df.copy()
    plot_df["label"] = plot_df["technology"].map(labels).fillna(plot_df["technology"])
    plot_df = plot_df.sort_values("probability_rank_1")
    fig, ax = plt.subplots(figsize=(8.5, 4.6), dpi=180)
    ax.barh(plot_df["label"], plot_df["probability_top_3"] * 100, color="#9ab6d8", label="Top 3")
    ax.barh(plot_df["label"], plot_df["probability_rank_1"] * 100, color="#2c5c9e", label="Rank 1")
    ax.set_title(title, loc="left", fontsize=12, fontweight="bold")
    ax.set_xlabel("Probability (%)")
    ax.set_xlim(0, 100)
    ax.legend(loc="lower right", fontsize=8)
    ax.grid(axis="x", color="#e5e7eb", linestyle=(0, (2, 4)))
    for spine in ("top", "right", "left"):
        ax.spines[spine].set_visible(False)
    ax.tick_params(axis="y", length=0, labelsize=8)
    ax.tick_params(axis="x", labelsize=8)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def save_scenario_chart(df: pd.DataFrame, sector: str, path: Path, title: str) -> None:
    sector_df = df.loc[df["sector"] == sector].copy()
    order = (
        sector_df.loc[sector_df["scenario"] == "Medium"]
        .sort_values("mean_specific", ascending=False)["technology"]
        .tolist()
    )
    pivot = sector_df.pivot(index="technology", columns="scenario", values="mean_specific").reindex(order)
    fig, ax = plt.subplots(figsize=(8.5, 4.6), dpi=180)
    y = np.arange(len(pivot.index))
    offsets = {"Low": -0.22, "Medium": 0.0, "High": 0.22}
    colors = {"Low": "#7aa6c2", "Medium": "#2c5c9e", "High": "#b85252"}
    for scenario in ["Low", "Medium", "High"]:
        ax.barh(y + offsets[scenario], pivot[scenario], height=0.2, label=scenario, color=colors[scenario])
    ax.axvline(0, color="#222222", linewidth=0.8)
    ax.set_yticks(y)
    ax.set_yticklabels([label_map(sector.lower()).get(t, t) for t in pivot.index], fontsize=8)
    ax.invert_yaxis()
    ax.set_title(title, loc="left", fontsize=12, fontweight="bold")
    ax.set_xlabel("Mean specific NPV")
    ax.legend(loc="lower right", fontsize=8)
    ax.grid(axis="x", color="#e5e7eb", linestyle=(0, (2, 4)))
    for spine in ("top", "right", "left"):
        ax.spines[spine].set_visible(False)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def write_report(
    electricity_summary: pd.DataFrame,
    cement_summary: pd.DataFrame,
    electricity_rank: pd.DataFrame,
    cement_rank: pd.DataFrame,
    co2_table: pd.DataFrame,
    discount_table: pd.DataFrame,
    macc_det: pd.DataFrame,
    macc_sim: pd.DataFrame,
) -> None:
    def table(df: pd.DataFrame, cols: list[str]) -> str:
        selected = df.loc[:, cols].copy()
        rendered_rows = []
        for _, row in selected.iterrows():
            rendered = []
            for value in row:
                if isinstance(value, float):
                    rendered.append(f"{value:.2f}")
                else:
                    rendered.append(str(value))
            rendered_rows.append(rendered)
        header = "| " + " | ".join(cols) + " |"
        divider = "| " + " | ".join("---" for _ in cols) + " |"
        body = ["| " + " | ".join(row) + " |" for row in rendered_rows]
        return "\n".join([header, divider, *body])

    lines = [
        "# Midterm Deep Dive Findings",
        "",
        "Generated from current source code with 100,000 Monte Carlo draws and random seed 42.",
        "",
        "## What we built",
        "",
        "- A two-sector probabilistic NPV model for electricity and cement.",
        "- Deterministic and Monte Carlo calculations for all modeled technologies.",
        "- Shared run-id logic so rankings compare technologies under common sampled market conditions.",
        "- Total and sector-normalized NPV outputs, ranking probabilities, sign counts, sensitivity heatmaps, scenario analysis, and cement MACC.",
        "",
        "## Current electricity findings",
        "",
        table(electricity_summary, ["technology", "mean_specific", "p05_specific", "p95_specific", "non_negative_share"]),
        "",
        table(electricity_rank, ["technology", "average_rank", "probability_rank_1", "probability_top_3"]),
        "",
        "## Current cement findings",
        "",
        table(cement_summary, ["technology", "mean_specific", "p05_specific", "p95_specific", "non_negative_share"]),
        "",
        table(cement_rank, ["technology", "average_rank", "probability_rank_1", "probability_top_3"]),
        "",
        "## Scenario findings",
        "",
        "CO2-price scenarios change fossil and directly emitting technologies most. Zero-direct-emission electricity technologies do not move with CO2 price.",
        "",
        "Discount-rate scenarios strongly affect capital-heavy technologies such as PV, wind, nuclear, cement CCS, and low-carbon cement routes.",
        "",
        "## Current cement MACC",
        "",
        table(macc_det, ["technology", "annual_abatement_tco2", "abatement_cost_eur_per_tco2"]),
        "",
        table(macc_sim, ["technology", "annual_abatement_tco2", "abatement_cost_eur_per_tco2", "abatement_cost_p05_eur_per_tco2", "abatement_cost_p95_eur_per_tco2"]),
        "",
    ]
    REPORT.write_text("\n".join(lines))


def add_table_cards(slide, rows, x, y, w, title, accent):
    body = "\n".join(rows)
    add_card(slide, title, body, x, y, w, Inches(2.1), accent, 15, 11)


def build_deck() -> tuple[Path, int]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    electricity = simulate_electricity_results(
        sample_size=SAMPLE_SIZE,
        random_seed=RANDOM_SEED,
        technologies=tuple(ELECTRICITY_TECHNOLOGY_LABELS),
    )
    cement = simulate_cement_results(
        sample_size=SAMPLE_SIZE,
        random_seed=RANDOM_SEED,
        technologies=tuple(CEMENT_TECHNOLOGY_LABELS),
        retrofit_bau_mode=DEFAULT_RETROFIT_BAU_MODE,
    )
    electricity_summary = summarize(
        electricity,
        ELECTRICITY_TECHNOLOGY_LABELS,
        "npv_eur_per_mwh",
    )
    cement_summary = summarize(cement, CEMENT_TECHNOLOGY_LABELS, "npv_eur_per_t")
    _, electricity_rank = calculate_electricity_npv_rankings_from_results(
        electricity,
        "Electricity",
        npv_scale="EUR/MWh",
    )
    _, cement_rank = calculate_cement_npv_rankings_from_results(
        cement,
        "Cement",
        npv_scale="EUR/t",
    )
    co2_table = pd.concat(
        [
            build_scenario_table("Electricity", electricity, ELECTRICITY_TECHNOLOGY_LABELS, "lifetime_output_mwh", "co2"),
            build_scenario_table("Cement", cement, CEMENT_TECHNOLOGY_LABELS, "lifetime_output_t", "co2"),
        ],
        ignore_index=True,
    )
    discount_table = pd.concat(
        [
            build_scenario_table("Electricity", electricity, ELECTRICITY_TECHNOLOGY_LABELS, "lifetime_output_mwh", "discount"),
            build_scenario_table("Cement", cement, CEMENT_TECHNOLOGY_LABELS, "lifetime_output_t", "discount"),
        ],
        ignore_index=True,
    )
    sensitivity = pd.concat(
        [
            standardized_sensitivity("electricity", 0.20, METRIC_SPECIFIC),
            standardized_sensitivity("cement", 0.20, METRIC_SPECIFIC),
        ],
        ignore_index=True,
    )
    macc_det = deterministic_cement_macc()
    macc_sim = simulated_cement_macc(sample_size=SAMPLE_SIZE, random_seed=RANDOM_SEED)

    save_bar_chart(electricity_summary, ASSET_DIR / "electricity_mean_specific.png", "Electricity mean NPV", "EUR/MWh")
    save_bar_chart(cement_summary, ASSET_DIR / "cement_mean_specific.png", "Cement mean NPV", "EUR/t")
    save_rank_chart(electricity_rank, ELECTRICITY_TECHNOLOGY_LABELS, ASSET_DIR / "electricity_rank_probability.png", "Electricity rank probabilities")
    save_rank_chart(cement_rank, CEMENT_TECHNOLOGY_LABELS, ASSET_DIR / "cement_rank_probability.png", "Cement rank probabilities")
    save_scenario_chart(co2_table, "Electricity", ASSET_DIR / "electricity_co2_scenarios.png", "Electricity CO2-price scenario")
    save_scenario_chart(co2_table, "Cement", ASSET_DIR / "cement_co2_scenarios.png", "Cement CO2-price scenario")
    save_scenario_chart(discount_table, "Electricity", ASSET_DIR / "electricity_discount_scenarios.png", "Electricity discount-rate scenario")
    save_scenario_chart(discount_table, "Cement", ASSET_DIR / "cement_discount_scenarios.png", "Cement discount-rate scenario")
    for sector in ("electricity", "cement"):
        fig = build_sensitivity_heatmap_figure(sensitivity, sector, 0.20, METRIC_SPECIFIC)
        fig.savefig(ASSET_DIR / f"{sector}_sensitivity_current.png", bbox_inches="tight")
        plt.close(fig)
    fig = plot_cement_macc(macc_det, title="Cement MACC - current deterministic")
    fig.savefig(ASSET_DIR / "cement_macc_current_deterministic.png", bbox_inches="tight")
    plt.close(fig)

    write_report(
        electricity_summary,
        cement_summary,
        electricity_rank,
        cement_rank,
        co2_table,
        discount_table,
        macc_det,
        macc_sim,
    )

    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    slides = []

    s = blank(prs)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = COLORS["light"]
    add_textbox(s, "Midterm Deep Dive", Inches(0.75), Inches(0.95), Inches(9.0), Inches(0.55), 36, True, "blue")
    add_textbox(s, "What we built, what we tested, and what we found", Inches(0.78), Inches(1.65), Inches(10.4), Inches(0.45), 22, True)
    add_textbox(s, "Current-source results | 100,000 Monte Carlo draws | seed 42", Inches(0.8), Inches(2.25), Inches(8.8), Inches(0.3), 13, False, "muted")
    add_card(s, "Research spine", "Research question -> experiment setup -> assumptions -> results -> discussion -> quality outlook.", Inches(0.9), Inches(4.1), Inches(5.5), Inches(1.25), "green")
    add_card(s, "Key improvement", "This version interprets sensitivity and scenario analysis instead of only showing figures.", Inches(6.8), Inches(4.1), Inches(5.5), Inches(1.25), "blue")
    slides.append(s)

    s = blank(prs)
    add_title(s, "Research Question", "What the thesis is trying to answer")
    add_card(s, "Main question", "Which emitting sectors should implement carbon capture, and which should investigate alternative decarbonization technologies?", Inches(1.0), Inches(1.55), Inches(11.3), Inches(1.15), "blue", 18, 15)
    add_bullets(s, [
        "The thesis is not asking whether CCS works technically; it asks where CCS is economically preferred under uncertainty.",
        "Carbon capture demand is therefore treated as a probability distribution generated by competition between technologies.",
        "The midterm work tests this framework on electricity and cement before extending to aluminium and chemicals.",
    ], Inches(1.05), Inches(3.25), Inches(10.9), Inches(2.0), 18)
    add_footer(s, 2)
    slides.append(s)

    s = blank(prs)
    add_title(s, "What We Built", "Project work completed so far")
    add_bullets(s, [
        "Electricity NPV model: deterministic and Monte Carlo calculations for fossil, CCS, nuclear, renewable, and biogas technologies.",
        "Cement NPV model: deterministic and Monte Carlo calculations for BAU, CCS, alternative production routes, and BAU-relative retrofits.",
        "Aligned Monte Carlo run IDs: each simulation ID represents one shared uncertain world for fair rankings.",
        "Comparison layer: total NPV, normalized NPV, rank probabilities, sign counts, sensitivity heatmaps, scenario analysis, and cement MACC.",
        "Documentation layer: README, handover guide, changelog, notebooks, figures, and reproducible CLI workflows.",
    ], Inches(0.9), Inches(1.55), Inches(11.2), Inches(4.2), 17)
    add_footer(s, 3)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Experiment Setup", "How the current research experiment works")
    add_metric(s, "1 TWh/y", "electricity output", Inches(0.9), Inches(1.55), Inches(2.4), "blue")
    add_metric(s, "1 Mt/y", "cement output", Inches(3.9), Inches(1.55), Inches(2.4), "green")
    add_metric(s, "80", "EUR/tCO2", Inches(6.9), Inches(1.55), Inches(2.4), "amber")
    add_metric(s, "8%", "discount rate", Inches(9.9), Inches(1.55), Inches(2.4), "blue")
    add_card(s, "Decision metric", "NPV is computed from annual revenue, CAPEX, fixed OPEX, variable OPEX, fuel/electricity cost, and carbon cost. Higher NPV is better.", Inches(0.9), Inches(3.15), Inches(5.75), Inches(1.35), "blue")
    add_card(s, "Uncertainty metric", "Rank probabilities estimate how often a technology wins or lands in the top three under sampled techno-economic futures.", Inches(6.9), Inches(3.15), Inches(5.45), Inches(1.35), "green")
    add_footer(s, 4)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Assumptions We Had To Make", "Important modeling choices")
    add_bullets(s, [
        "Market and policy assumptions: carbon price, discount rate, fuel prices, electricity price, and cement price.",
        "Technology assumptions: CAPEX, OPEX, energy intensity, emissions intensity, full-load hours, and lifetime.",
        "Cement retrofits are modeled relative to BAU; CCS strongly reduces emissions but can increase energy use.",
        "Alternative fuels were corrected: fuel costs now use a 25-60% alternative-fuel share and blended fuel price.",
        "Fuel use and emissions are currently sampled independently; this is a simplifying assumption that should be tested in future work.",
    ], Inches(0.9), Inches(1.55), Inches(11.1), Inches(4.1), 17)
    add_footer(s, 5)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Result 1: Electricity Baseline", "Current model says electricity CCS demand is weak")
    add_picture_fit(s, ASSET_DIR / "electricity_mean_specific.png", Inches(0.65), Inches(1.4), Inches(6.15), Inches(4.75))
    add_picture_fit(s, ASSET_DIR / "electricity_rank_probability.png", Inches(6.95), Inches(1.4), Inches(5.75), Inches(4.75))
    add_footer(s, 6)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Electricity Interpretation", "What happened in the electricity model")
    add_bullets(s, [
        f"PV is the strongest technology in mean NPV: {electricity_summary.iloc[0]['mean_specific']:.2f} EUR/MWh.",
        f"PV and onshore wind both have almost always non-negative NPV; offshore wind is non-negative in {electricity_summary.loc[electricity_summary.technology == 'wind_offshore', 'non_negative_share'].iloc[0] * 100:.0f}% of draws.",
        "CCGT is the best fossil option, but remains negative on average and only occasionally reaches positive NPV.",
        "CCGT+CCS and hard coal+CCS reduce emissions exposure but do not overcome added cost and energy penalties.",
        "Nuclear remains negative on average because high CAPEX dominates even with high full-load hours and long lifetime.",
    ], Inches(0.9), Inches(1.55), Inches(11.1), Inches(4.0), 18)
    add_footer(s, 7)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Result 2: Cement Baseline", "Cement is the stronger CCS-demand candidate")
    add_picture_fit(s, ASSET_DIR / "cement_mean_specific.png", Inches(0.65), Inches(1.4), Inches(6.15), Inches(4.75))
    add_picture_fit(s, ASSET_DIR / "cement_rank_probability.png", Inches(6.95), Inches(1.4), Inches(5.75), Inches(4.75))
    add_footer(s, 8)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Cement Interpretation", "What happened in the cement model")
    add_bullets(s, [
        f"CCS has the highest mean cement NPV in the current 100,000-draw run: {cement_summary.loc[cement_summary.technology == 'ccs', 'mean_specific'].iloc[0]:.2f} EUR/t.",
        f"CCS also has the highest rank-1 probability: {cement_rank.loc[cement_rank.technology == 'ccs', 'probability_rank_1'].iloc[0] * 100:.1f}%.",
        "Process heat integration, waste heat recovery, clinker substitution, and alternative fuels cluster close to CCS, so CCS is competitive rather than uncontested.",
        "Electrification and electrolysis are strongly negative because electricity demand dominates the economics under current electricity-price assumptions.",
        "All conventional and retrofit cement options are non-negative in the sampled baseline; the hard problem is ranking and abatement depth, not basic viability.",
    ], Inches(0.9), Inches(1.55), Inches(11.1), Inches(4.0), 18)
    add_footer(s, 9)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Result 3: Scenario Analysis", "CO2 price and discount rate tell different stories")
    add_picture_fit(s, ASSET_DIR / "electricity_co2_scenarios.png", Inches(0.55), Inches(1.35), Inches(6.05), Inches(2.45))
    add_picture_fit(s, ASSET_DIR / "cement_co2_scenarios.png", Inches(6.85), Inches(1.35), Inches(6.05), Inches(2.45))
    add_picture_fit(s, ASSET_DIR / "electricity_discount_scenarios.png", Inches(0.55), Inches(4.15), Inches(6.05), Inches(2.45))
    add_picture_fit(s, ASSET_DIR / "cement_discount_scenarios.png", Inches(6.85), Inches(4.15), Inches(6.05), Inches(2.45))
    add_footer(s, 10)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Scenario Interpretation", "What we learned from the scenario notebook")
    add_bullets(s, [
        "CO2 price mostly affects technologies with direct emissions. PV, wind, nuclear, and biogas do not move in the direct-emissions-only setup.",
        "Higher CO2 price makes unabated coal and gas worse, so CCS narrows the fossil gap, but does not make fossil CCS beat renewables in electricity.",
        "In cement, higher CO2 price lowers all directly emitting routes; CCS moves less because residual emissions are small.",
        "Lower discount rates strongly improve capital-heavy low-carbon technologies, especially PV, wind, nuclear, and cement CCS.",
        "Even at 4% discount rate, nuclear remains close to break-even but still negative on average in the current 100,000-draw run, so its problem is not only financing; CAPEX remains too high relative to revenue.",
    ], Inches(0.9), Inches(1.55), Inches(11.2), Inches(4.2), 17)
    add_footer(s, 11)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Result 4: Sensitivity Analysis", "The heatmaps show why scenarios move the way they do")
    add_picture_fit(s, ASSET_DIR / "electricity_sensitivity_current.png", Inches(0.55), Inches(1.35), Inches(6.05), Inches(4.85))
    add_picture_fit(s, ASSET_DIR / "cement_sensitivity_current.png", Inches(6.85), Inches(1.35), Inches(6.05), Inches(4.85))
    add_footer(s, 12)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Sensitivity Interpretation", "Main drivers by technology family")
    add_table_cards(s, [
        "Renewables: full-load hours, CAPEX, discount rate.",
        "Fossil electricity: fuel price/use, emissions, carbon price.",
        "Nuclear: full-load hours, CAPEX, lifetime.",
    ], Inches(0.8), Inches(1.55), Inches(5.7), "Electricity", "blue")
    add_table_cards(s, [
        "Conventional/retrofit cement: emissions, carbon price, discount rate.",
        "CCS cement: discount rate plus electricity penalty.",
        "Electrification/electrolysis: electricity use and electricity price.",
    ], Inches(6.85), Inches(1.55), Inches(5.7), "Cement", "green")
    add_bullets(s, [
        "This is why the scenario analysis should be discussed, not merely shown: it explains which uncertain assumptions can change the CCS-demand conclusion.",
        "Sensitivity results also point to quality improvements: better electricity-price scenarios, better emissions modeling, and clearer source quality for technology assumptions.",
    ], Inches(1.0), Inches(4.25), Inches(11.1), Inches(1.3), 17)
    add_footer(s, 13)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Result 5: Cement MACC", "NPV ranking alone misses abatement depth")
    add_picture_fit(s, ASSET_DIR / "cement_macc_current_deterministic.png", Inches(0.55), Inches(1.35), Inches(7.1), Inches(4.75))
    add_card(s, "MACC finding", f"CCS avoids about {macc_det.loc[macc_det.technology == 'ccs', 'annual_abatement_tco2'].iloc[0] / 1e6:.2f} MtCO2/y at {macc_det.loc[macc_det.technology == 'ccs', 'abatement_cost_eur_per_tco2'].iloc[0]:.1f} EUR/tCO2 in the deterministic current model.", Inches(8.0), Inches(1.6), Inches(4.35), Inches(1.45), "green", 15, 12)
    add_card(s, "Interpretation", "Efficiency and heat-integration measures can be cheap or cost-saving, but their abatement depth is small. CCS is more expensive than small retrofits, but its abatement depth is much larger.", Inches(8.0), Inches(3.45), Inches(4.35), Inches(1.8), "blue", 15, 12)
    add_footer(s, 14)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Discussion", "Current answer to the thesis question")
    add_bullets(s, [
        "Electricity: CCS demand looks weak under current assumptions because non-capture renewables are economically dominant.",
        "Cement: CCS demand looks conditional but significant because CCS is often economically competitive and provides deep direct-emissions abatement.",
        "The thesis mechanism is visible: CCS demand depends on sector-specific alternatives, not only on the cost of CCS itself.",
        "The model should eventually classify sectors into robust CCS demand, conditional CCS demand, weak CCS demand, and unresolved cases.",
    ], Inches(0.9), Inches(1.55), Inches(11.2), Inches(3.6), 18)
    add_footer(s, 15)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Quality-Level Improvements", "How the research should get better")
    add_bullets(s, [
        "Strengthen parameter sourcing and uncertainty ranges, especially for emerging industrial technologies.",
        "Separate technology-development uncertainty from market and policy uncertainty more explicitly.",
        "Replace weak independent assumptions with mechanistic links where possible: fuel use, emissions, capture rate, and energy penalty.",
        "Add incremental comparisons: CCS versus unabated counterpart in electricity, and retrofits versus BAU in cement.",
        "Add abatement-depth interpretation alongside NPV ranking so small cheap measures are not overinterpreted as complete CCS substitutes.",
    ], Inches(0.9), Inches(1.55), Inches(11.1), Inches(4.0), 17)
    add_footer(s, 16)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Outlook", "How the project proceeds scientifically")
    add_card(s, "Aluminium", "Add aluminium industry pathways and compare CCS with process integration, electrification, inert anodes, recycling, and fuel switching where relevant.", Inches(0.8), Inches(1.55), Inches(5.7), Inches(1.55), "blue")
    add_card(s, "Chemical industries", "Add chemical-sector cases where feedstock, process emissions, hydrogen, electrification, and CCU/CCS alternatives differ by product.", Inches(6.85), Inches(1.55), Inches(5.7), Inches(1.55), "green")
    add_card(s, "Cross-sector conclusion", "Move from technology rankings to sector-level CCS demand categories and identify the conditions that make each category robust.", Inches(0.8), Inches(3.65), Inches(11.75), Inches(1.2), "amber")
    add_footer(s, 17)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Technical Steps TODO", "Implementation items, separate from research outlook")
    add_bullets(s, [
        "Regenerate official tracked cement figures after the alternative-fuels correction.",
        "Add tests for finance identities, unit conversions, sign counts, and ranking consistency.",
        "Write reusable sector templates for aluminium and chemicals.",
        "Archive parameter sources next to assumptions.",
        "Record sample size, seed, scenario values, and model version for every exported figure.",
    ], Inches(0.9), Inches(1.55), Inches(11.2), Inches(3.4), 18)
    add_footer(s, 18)
    slides.append(s)

    prs.save(OUT)
    return OUT, len(slides)


if __name__ == "__main__":
    path, count = build_deck()
    print(f"Wrote {path} with {count} slides")
    print(f"Wrote {REPORT}")
