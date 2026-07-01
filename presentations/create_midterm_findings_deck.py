from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentations" / "midterm_findings_v1.pptx"

FIG = ROOT / "figures"
DATA = ROOT / "data" / "processed"

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

COLORS = {
    "ink": RGBColor(31, 37, 48),
    "muted": RGBColor(89, 99, 115),
    "line": RGBColor(210, 216, 224),
    "blue": RGBColor(44, 92, 158),
    "green": RGBColor(40, 130, 92),
    "red": RGBColor(176, 72, 72),
    "amber": RGBColor(192, 132, 46),
    "light": RGBColor(246, 248, 251),
    "white": RGBColor(255, 255, 255),
}


def add_textbox(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    font_size=20,
    bold=False,
    color="ink",
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_title(slide, title: str, subtitle: str | None = None):
    add_textbox(slide, title, Inches(0.55), Inches(0.28), Inches(12.1), Inches(0.55), 28, True)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.58), Inches(0.82), Inches(11.8), Inches(0.35), 12, False, "muted")
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(1.2),
        Inches(12.2),
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def add_footer(slide, n: int):
    add_textbox(
        slide,
        f"Midterm findings | {n}",
        Inches(11.3),
        Inches(7.05),
        Inches(1.55),
        Inches(0.22),
        8,
        False,
        "muted",
        PP_ALIGN.RIGHT,
    )


def add_bullets(slide, bullets, x, y, w, h, font_size=17, color="ink", gap=0.15):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(gap * 12)
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(font_size)
        run.font.color.rgb = COLORS[color]
    return box


def add_card(slide, title, body, x, y, w, h, accent="blue", title_size=15, body_size=12):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    shape.adjustments[0] = 0.08
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[accent]
    bar.line.fill.background()
    add_textbox(slide, title, x + Inches(0.18), y + Inches(0.12), w - Inches(0.28), Inches(0.28), title_size, True, accent)
    add_textbox(slide, body, x + Inches(0.18), y + Inches(0.47), w - Inches(0.28), h - Inches(0.55), body_size, False, "ink")


def add_picture_fit(slide, image_path: Path, x, y, w, h):
    from PIL import Image

    with Image.open(image_path) as im:
        iw, ih = im.size
    image_ratio = iw / ih
    box_ratio = w / h
    if image_ratio > box_ratio:
        pic_w = w
        pic_h = w / image_ratio
        pic_x = x
        pic_y = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * image_ratio
        pic_x = x + (w - pic_w) / 2
        pic_y = y
    return slide.shapes.add_picture(str(image_path), pic_x, pic_y, pic_w, pic_h)


def add_metric(slide, value, label, x, y, w, accent="blue"):
    add_textbox(slide, value, x, y, w, Inches(0.36), 23, True, accent, PP_ALIGN.CENTER)
    add_textbox(slide, label, x, y + Inches(0.36), w, Inches(0.42), 9, False, "muted", PP_ALIGN.CENTER)


def read_csv(path: Path):
    with path.open(newline="") as f:
        return list(csv.DictReader(f))


def fmt_pct(value):
    return f"{float(value) * 100:.1f}%"


def fmt_rank(value):
    return f"{float(value):.2f}"


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_deck():
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H

    elec_rank = read_csv(DATA / "2026-06-23-NPV_Ranking_Electricity_summary.csv")
    cement_rank = read_csv(DATA / "2026-06-23-NPV_Ranking_Cement_summary.csv")
    macc_det = read_csv(DATA / "2026-06-26-Cement_MACC_Deterministic.csv")
    macc_sim = read_csv(DATA / "2026-06-26-Cement_MACC_Simulated.csv")

    slides = []

    s = blank(prs)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = COLORS["light"]
    add_textbox(s, "Midterm Findings", Inches(0.75), Inches(1.05), Inches(6.3), Inches(0.55), 36, True, "blue")
    add_textbox(
        s,
        "How uncertainty of technology development shapes carbon capture demand uncertainty",
        Inches(0.78),
        Inches(1.72),
        Inches(8.7),
        Inches(0.5),
        21,
        True,
    )
    add_textbox(
        s,
        "Finn | Master thesis project | ETH Zurich | July 2026",
        Inches(0.8),
        Inches(2.35),
        Inches(8.2),
        Inches(0.32),
        13,
        False,
        "muted",
    )
    add_card(
        s,
        "Midterm thesis question",
        "Under techno-economic uncertainty, when is carbon capture the most economically viable option, and when do non-capture alternatives dominate?",
        Inches(0.85),
        Inches(4.15),
        Inches(5.2),
        Inches(1.35),
        "green",
        15,
        13,
    )
    add_card(
        s,
        "Current implemented scope",
        "Two mature sectors: electricity and cement. Both use deterministic calculations, Monte Carlo rankings, sensitivity analysis, and tracked figures.",
        Inches(6.4),
        Inches(4.15),
        Inches(5.9),
        Inches(1.35),
        "blue",
        15,
        13,
    )
    slides.append(s)

    s = blank(prs)
    add_title(s, "What The Thesis Asks", "Reframed from the project proposal")
    add_bullets(
        s,
        [
            "CCS is important for net-zero, but deployment can be reduced if plants choose alternative decarbonization technologies.",
            "The thesis therefore studies CCS demand as a competitive investment problem, not as a standalone technology assessment.",
            "The core output is probabilistic: how often does CCS become the economically preferred option under uncertain technology development?",
        ],
        Inches(0.75),
        Inches(1.65),
        Inches(6.2),
        Inches(2.1),
        18,
    )
    add_card(
        s,
        "Research question",
        "Which emitting sectors should implement carbon capture, and which should investigate alternative decarbonization technologies?",
        Inches(7.35),
        Inches(1.65),
        Inches(4.9),
        Inches(1.28),
        "blue",
        15,
        13,
    )
    add_card(
        s,
        "Midterm answer format",
        "For each sector: compare CCS against non-capture options using NPV distributions, rank probabilities, sensitivity heatmaps, and marginal abatement costs.",
        Inches(7.35),
        Inches(3.25),
        Inches(4.9),
        Inches(1.45),
        "green",
        15,
        13,
    )
    add_footer(s, 2)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Current Scope", "Implemented applications and comparison technologies")
    add_card(
        s,
        "Electricity",
        "Hard coal, hard coal + CCS, CCGT, CCGT + CCS, nuclear, offshore wind, onshore wind, PV, biogas.",
        Inches(0.75),
        Inches(1.55),
        Inches(5.8),
        Inches(1.45),
        "blue",
    )
    add_card(
        s,
        "Cement",
        "BAU, CCS, clinker substitution, alternative fuels, efficiency improvement, waste heat recovery, process heat integration, electrification, electrolysis.",
        Inches(6.85),
        Inches(1.55),
        Inches(5.8),
        Inches(1.45),
        "green",
    )
    add_bullets(
        s,
        [
            "Both sectors are normalized to a common annual output before comparison.",
            "The model compares technology choices within a sector, not electricity versus cement directly.",
            "Steel, metals, and chemicals remain thesis-scope candidates, but are not yet mature midterm result sectors.",
        ],
        Inches(1.0),
        Inches(3.6),
        Inches(11.3),
        Inches(1.9),
        17,
    )
    add_footer(s, 3)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Model Logic", "From assumptions to CCS demand insight")
    steps = [
        ("1. Inputs", "Technology assumptions, market prices, carbon price, discount rate.", "blue"),
        ("2. Sampling", "Monte Carlo draws techno-economic uncertainty for each simulation world.", "green"),
        ("3. Cash Flow", "Annual revenue minus CAPEX, OPEX, fuel, electricity, and carbon costs.", "amber"),
        ("4. NPV", "Discount future annual cash flows over the technology lifetime.", "blue"),
        ("5. Decision", "Rank technologies and estimate the probability that CCS is preferred.", "green"),
    ]
    x = Inches(0.65)
    for title, body, color in steps:
        add_card(s, title, body, x, Inches(2.0), Inches(2.35), Inches(1.6), color, 14, 11)
        x += Inches(2.45)
    add_bullets(
        s,
        [
            "The same sampled market conditions are reused across comparable technologies where appropriate.",
            "Higher NPV is better; NPV = 0 is treated as non-negative.",
            "Outputs are reported as total NPV and normalized NPV: EUR/MWh for electricity and EUR/t for cement.",
        ],
        Inches(0.9),
        Inches(4.35),
        Inches(11.5),
        Inches(1.2),
        16,
    )
    add_footer(s, 4)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Core Assumptions", "Common assumptions used in the current model")
    metrics = [
        ("80", "EUR/tCO2 carbon price", "blue"),
        ("8%", "discount rate", "green"),
        ("100k", "Monte Carlo draws", "amber"),
        ("42", "default random seed", "blue"),
    ]
    for i, (v, lab, c) in enumerate(metrics):
        add_metric(s, v, lab, Inches(0.9 + i * 3.05), Inches(1.55), Inches(2.45), c)
    add_card(
        s,
        "Electricity normalization",
        "All electricity technologies are sized to produce 1,000,000 MWh/year. Full-load hours determine installed capacity.",
        Inches(0.8),
        Inches(3.0),
        Inches(5.8),
        Inches(1.25),
        "blue",
    )
    add_card(
        s,
        "Cement normalization",
        "All cement technologies are compared at 1,000,000 t/year. Cement retrofits modify a BAU baseline.",
        Inches(6.85),
        Inches(3.0),
        Inches(5.8),
        Inches(1.25),
        "green",
    )
    add_card(
        s,
        "Current revenue assumptions",
        "Electricity price is fixed at 94.07 EUR/MWh for electricity-sector revenue. Cement price is fixed at 150 EUR/t for cement-sector revenue.",
        Inches(0.8),
        Inches(4.65),
        Inches(11.85),
        Inches(0.95),
        "amber",
        14,
        12,
    )
    add_footer(s, 5)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Shared Market Inputs", "Uncertain energy prices and fixed fuel assumptions")
    add_card(s, "Gas price", "Scaled beta: 12.5 / 39.3 / 89.7 EUR/MWh_th", Inches(0.8), Inches(1.55), Inches(3.8), Inches(1.05), "blue")
    add_card(s, "Coal price", "Scaled beta: 8.14 / 12.11 / 24.17 EUR/MWh_th", Inches(4.8), Inches(1.55), Inches(3.8), Inches(1.05), "blue")
    add_card(s, "Electricity price", "Scaled beta: 74.8 / 183.7 / 255.2 EUR/MWh", Inches(8.8), Inches(1.55), Inches(3.8), Inches(1.05), "blue")
    add_card(s, "Biofuel price", "Uniform: 5.4-32.4 EUR/MWh_th", Inches(0.8), Inches(3.0), Inches(3.8), Inches(1.05), "green")
    add_card(s, "Nuclear fuel", "Fixed: 2.8 EUR/MWh_th", Inches(4.8), Inches(3.0), Inches(3.8), Inches(1.05), "green")
    add_card(s, "Biogas fuel", "Fixed: 87.5 EUR/MWh_th", Inches(8.8), Inches(3.0), Inches(3.8), Inches(1.05), "green")
    add_bullets(
        s,
        [
            "Energy-price assumptions are intentionally separated from technology-specific assumptions.",
            "This keeps technology development uncertainty and market/policy uncertainty traceable.",
        ],
        Inches(1.0),
        Inches(4.75),
        Inches(11.2),
        Inches(0.9),
        16,
    )
    add_footer(s, 6)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Technology Parameters", "What is sampled or fixed for each option")
    add_bullets(
        s,
        [
            "CAPEX: upfront investment cost.",
            "Fixed OPEX: annual capacity- or output-linked fixed operating cost.",
            "Variable OPEX: non-fuel variable operating cost.",
            "Fuel and electricity intensity: translate energy requirements into annual costs.",
            "Direct emissions: translate residual emissions into carbon cost.",
            "Lifetime and full-load hours: determine discount horizon and installed capacity.",
        ],
        Inches(0.9),
        Inches(1.55),
        Inches(5.8),
        Inches(3.5),
        16,
    )
    add_card(
        s,
        "Cement retrofit detail",
        "Retrofits are represented as changes relative to BAU. CCS, for example, lowers emissions strongly but can increase fuel and electricity use.",
        Inches(7.2),
        Inches(1.6),
        Inches(4.95),
        Inches(1.35),
        "green",
    )
    add_card(
        s,
        "Recent assumption correction",
        "Alternative fuels now use a 25-60% alternative-fuel share and a blended fuel price, rather than assuming 100% alternative fuel use.",
        Inches(7.2),
        Inches(3.35),
        Inches(4.95),
        Inches(1.35),
        "amber",
    )
    add_footer(s, 7)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Electricity Results", "Mean NPV and rank probabilities")
    add_picture_fit(s, FIG / "2026-06-23-Average_NPV_Rank_Electricity.png", Inches(0.65), Inches(1.45), Inches(7.15), Inches(4.7))
    top = {r["technology"]: r for r in elec_rank}
    add_card(
        s,
        "Dominant options",
        f"PV: average rank {fmt_rank(top['pv']['average_rank'])}, rank-1 probability {fmt_pct(top['pv']['probability_rank_1'])}.\n"
        f"Onshore wind: average rank {fmt_rank(top['wind_onshore']['average_rank'])}, top-3 probability {fmt_pct(top['wind_onshore']['probability_top_3'])}.\n"
        f"Offshore wind: average rank {fmt_rank(top['wind_offshore']['average_rank'])}, top-3 probability {fmt_pct(top['wind_offshore']['probability_top_3'])}.",
        Inches(8.05),
        Inches(1.55),
        Inches(4.55),
        Inches(2.05),
        "green",
        15,
        12,
    )
    add_card(
        s,
        "CCS implication",
        f"CCGT+CCS and hard coal+CCS do not reach top-3 probability above {fmt_pct(top['ccgt_ccs']['probability_top_3'])} and {fmt_pct(top['hard_coal_ccs']['probability_top_3'])}, respectively.",
        Inches(8.05),
        Inches(3.85),
        Inches(4.55),
        Inches(1.4),
        "blue",
        15,
        12,
    )
    add_footer(s, 8)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Electricity Interpretation", "What the ranking says about CCS demand")
    add_bullets(
        s,
        [
            "In the current model, electricity-sector CCS demand is weak because renewable options dominate the economic ranking.",
            "CCS reduces carbon costs for fossil generation, but added CAPEX, OPEX, and fuel consumption keep CCS behind PV and wind.",
            "Nuclear is disadvantaged by high upfront CAPEX despite high full-load hours and a 45-year lifetime.",
            "Biogas is penalized by high fixed fuel cost.",
        ],
        Inches(0.9),
        Inches(1.6),
        Inches(6.1),
        Inches(3.3),
        17,
    )
    add_picture_fit(s, FIG / "2026-06-23-Mean_NPV_per_MWh_Electricity.png", Inches(7.2), Inches(1.45), Inches(5.2), Inches(4.7))
    add_footer(s, 9)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Cement Results", "A stronger but uncertain case for CCS")
    add_picture_fit(s, FIG / "2026-06-23-Average_NPV_Rank_Cement.png", Inches(0.65), Inches(1.45), Inches(7.15), Inches(4.7))
    ctop = {r["technology"]: r for r in cement_rank}
    add_card(
        s,
        "CCS probability signal",
        f"CCS has rank-1 probability {fmt_pct(ctop['ccs']['probability_rank_1'])}, the highest in cement. But its average rank is {fmt_rank(ctop['ccs']['average_rank'])}, showing high volatility.",
        Inches(8.05),
        Inches(1.55),
        Inches(4.55),
        Inches(1.55),
        "green",
        15,
        12,
    )
    add_card(
        s,
        "Competing non-capture options",
        f"Process heat integration, waste heat recovery, and clinker substitution all have top-3 probabilities above {fmt_pct(ctop['clinker_substitution']['probability_top_3'])}.",
        Inches(8.05),
        Inches(3.35),
        Inches(4.55),
        Inches(1.45),
        "blue",
        15,
        12,
    )
    add_footer(s, 10)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Cement Interpretation", "CCS competes with low-cost retrofits")
    add_bullets(
        s,
        [
            "Cement is currently the stronger CCS-demand candidate because CCS provides large direct emissions reduction.",
            "However, non-capture retrofits can remain economically attractive because they add little CAPEX and reduce energy or emissions intensity.",
            "Electrification and electrolysis are currently unattractive because electricity demand is high under the present electricity-price assumption.",
            "Therefore, cement CCS demand depends on the trade-off between abatement depth and retrofit cost.",
        ],
        Inches(0.9),
        Inches(1.55),
        Inches(6.15),
        Inches(3.4),
        17,
    )
    add_picture_fit(s, FIG / "2026-06-23-Mean_NPV_per_t_Cement.png", Inches(7.25), Inches(1.45), Inches(5.15), Inches(4.7))
    add_footer(s, 11)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Sensitivity: Cement", "Standardized +/-20% changes in model inputs")
    add_picture_fit(s, FIG / "2026-06-24-Sensitivity_Heatmap_Standardized_Cement.png", Inches(0.55), Inches(1.38), Inches(8.1), Inches(4.85))
    add_card(
        s,
        "Main insight",
        "Electricity-intensive cement routes are highly sensitive to electricity price and electricity use. Conventional and retrofit routes are more sensitive to emissions, carbon price, and financial assumptions.",
        Inches(8.9),
        Inches(1.55),
        Inches(3.55),
        Inches(1.95),
        "green",
        15,
        12,
    )
    add_card(
        s,
        "Interpretation for CCS",
        "CCS attractiveness is coupled to carbon price, residual emissions, energy penalty, and discount rate.",
        Inches(8.9),
        Inches(3.85),
        Inches(3.55),
        Inches(1.25),
        "blue",
        15,
        12,
    )
    add_footer(s, 12)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Sensitivity: Electricity", "Standardized +/-20% changes in model inputs")
    add_picture_fit(s, FIG / "2026-06-24-Sensitivity_Heatmap_Standardized_Electricity.png", Inches(0.55), Inches(1.38), Inches(8.1), Inches(4.85))
    add_card(
        s,
        "Main insight",
        "Fossil technologies are exposed to fuel price, fuel use, emissions, and carbon price. Renewable technologies are mainly exposed to investment cost, fixed OPEX, and full-load hours.",
        Inches(8.9),
        Inches(1.55),
        Inches(3.55),
        Inches(1.95),
        "blue",
        15,
        12,
    )
    add_card(
        s,
        "Interpretation for CCS",
        "CCS can reduce carbon-price exposure but adds cost and efficiency penalties, so it does not become dominant in electricity under the current setup.",
        Inches(8.9),
        Inches(3.85),
        Inches(3.55),
        Inches(1.45),
        "green",
        15,
        12,
    )
    add_footer(s, 13)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Cement MACC", "Abatement cost connects economics with emissions reduction")
    add_picture_fit(s, FIG / "2026-06-26-Cement_MACC_Deterministic.png", Inches(0.55), Inches(1.38), Inches(7.7), Inches(4.6))
    det = {r["technology"]: r for r in macc_det}
    sim = {r["technology"]: r for r in macc_sim}
    add_card(
        s,
        "Deterministic MACC",
        f"Efficiency improvement: {float(det['efficiency_improvement']['abatement_cost_eur_per_tco2']):.1f} EUR/tCO2.\n"
        f"Process heat integration: {float(det['process_heat_integration']['abatement_cost_eur_per_tco2']):.1f} EUR/tCO2.\n"
        f"CCS: {float(det['ccs']['abatement_cost_eur_per_tco2']):.1f} EUR/tCO2.",
        Inches(8.55),
        Inches(1.55),
        Inches(3.9),
        Inches(1.85),
        "green",
        15,
        12,
    )
    add_card(
        s,
        "Simulated CCS range",
        f"CCS median: {float(sim['ccs']['abatement_cost_median_eur_per_tco2']):.1f} EUR/tCO2.\n"
        f"p05-p95: {float(sim['ccs']['abatement_cost_p05_eur_per_tco2']):.1f} to {float(sim['ccs']['abatement_cost_p95_eur_per_tco2']):.1f} EUR/tCO2.",
        Inches(8.55),
        Inches(3.75),
        Inches(3.9),
        Inches(1.45),
        "blue",
        15,
        12,
    )
    add_footer(s, 14)
    slides.append(s)

    s = blank(prs)
    add_title(s, "What This Means For CCS Demand", "Midterm synthesis")
    add_card(
        s,
        "Electricity",
        "Low expected CCS demand in the current setup. Renewables dominate the NPV ranking; CCS is not frequently the preferred electricity option.",
        Inches(0.9),
        Inches(1.65),
        Inches(5.65),
        Inches(1.5),
        "blue",
    )
    add_card(
        s,
        "Cement",
        "Stronger CCS-demand signal. CCS often ranks first and provides high abatement, but low-cost retrofits remain strong competitors.",
        Inches(6.9),
        Inches(1.65),
        Inches(5.65),
        Inches(1.5),
        "green",
    )
    add_bullets(
        s,
        [
            "The thesis question is sector-specific: CCS is not universally dominant or universally unattractive.",
            "A useful demand indicator is the probability that CCS ranks first or in the top three, not only the mean NPV.",
            "Abatement depth matters: small retrofits can be cheap, while CCS can deliver larger emissions reductions.",
        ],
        Inches(1.0),
        Inches(4.0),
        Inches(11.2),
        Inches(1.35),
        17,
    )
    add_footer(s, 15)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Current Limitations", "What should be caveated in the midterm")
    add_bullets(
        s,
        [
            "Only electricity and cement are mature enough for results; other proposed sectors are not yet modeled.",
            "There is no formal automated test suite yet; validation uses compilation, smoke runs, and notebook execution.",
            "Some physical correlations are simplified, especially between fuel use and emissions.",
            "Results depend strongly on assumed product prices, carbon price, energy prices, and discount rate.",
            "Stored June 23 cement summary figures predate the June 29 alternative-fuels correction and should be regenerated before final claims.",
        ],
        Inches(0.9),
        Inches(1.55),
        Inches(11.3),
        Inches(3.8),
        17,
    )
    add_footer(s, 16)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Next Steps", "From midterm results to final thesis")
    add_card(
        s,
        "Regenerate",
        "Regenerate cement summary outputs after the alternative-fuels correction.",
        Inches(0.8),
        Inches(1.6),
        Inches(3.7),
        Inches(1.1),
        "amber",
    )
    add_card(
        s,
        "Extend",
        "Add further emitting sectors if needed: steel/metals/chemicals.",
        Inches(4.8),
        Inches(1.6),
        Inches(3.7),
        Inches(1.1),
        "blue",
    )
    add_card(
        s,
        "Improve uncertainty",
        "Add correlation checks and stronger unit-identity tests.",
        Inches(8.8),
        Inches(1.6),
        Inches(3.7),
        Inches(1.1),
        "green",
    )
    add_card(
        s,
        "Scenario analysis",
        "Use CO2 price and discount-rate scenarios to test when CCS becomes robust.",
        Inches(0.8),
        Inches(3.2),
        Inches(3.7),
        Inches(1.1),
        "blue",
    )
    add_card(
        s,
        "Final narrative",
        "Translate model outputs into CCS demand categories by sector and application.",
        Inches(4.8),
        Inches(3.2),
        Inches(3.7),
        Inches(1.1),
        "green",
    )
    add_card(
        s,
        "Presentation polish",
        "Update deck figures after reruns and add source references for final defense.",
        Inches(8.8),
        Inches(3.2),
        Inches(3.7),
        Inches(1.1),
        "amber",
    )
    add_footer(s, 17)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Appendix: Deterministic Baselines", "Representative calculations for checking intuition")
    add_picture_fit(s, FIG / "2026-06-23-Deterministic_NPV_Electricity.png", Inches(0.75), Inches(1.55), Inches(5.7), Inches(4.7))
    add_picture_fit(s, FIG / "2026-06-23-Deterministic_NPV_Cement.png", Inches(6.9), Inches(1.55), Inches(5.7), Inches(4.7))
    add_footer(s, 18)
    slides.append(s)

    s = blank(prs)
    add_title(s, "Appendix: Simulated Cement MACC", "Uncertainty around abatement costs")
    add_picture_fit(s, FIG / "2026-06-26-Cement_MACC_Simulated.png", Inches(0.75), Inches(1.45), Inches(11.8), Inches(4.8))
    add_footer(s, 19)
    slides.append(s)

    prs.save(OUT)
    return OUT, len(prs.slides)


if __name__ == "__main__":
    path, count = build_deck()
    print(f"Wrote {path} with {count} slides")
